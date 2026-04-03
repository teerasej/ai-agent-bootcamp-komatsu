#!/usr/bin/env python3
"""
Generate a 5-slide presentation for D3-M7: PV Agent — Integrate Agent into a Web Chat UI (C#).
Run: python generate_webapp_slides.py
Output: webapp_slides.pptx in the same directory.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── colour palette ───────────────────────────────────────────────
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE     = RGBColor(0xF5, 0xF5, 0xF5)
DARK_BG       = RGBColor(0x1E, 0x1E, 0x2E)
ACCENT_BLUE   = RGBColor(0x00, 0x78, 0xD4)
ACCENT_TEAL   = RGBColor(0x00, 0xB4, 0xA0)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
LIGHT_GRAY    = RGBColor(0xE0, 0xE0, 0xE0)
MID_GRAY      = RGBColor(0x99, 0x99, 0x99)
DARK_TEXT     = RGBColor(0x22, 0x22, 0x22)
CODE_BG       = RGBColor(0x2D, 0x2D, 0x2D)
SECTION_BG    = RGBColor(0x00, 0x3E, 0x6B)
LETSDOIT_BG   = RGBColor(0x00, 0x5A, 0x9E)
TRACK1_CLR    = RGBColor(0x00, 0x78, 0xD4)
COSMOS_CLR    = RGBColor(0x6B, 0x2F, 0xA0)

# ── slide dimensions (widescreen 16:9) ───────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H


# ═══════════════════════════════════════════════════════════════════
#  Helpers (mirrors generate_presentation.py)
# ═══════════════════════════════════════════════════════════════════

def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, fill_color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullets(slide, left, top, width, height, items, font_size=16,
                 color=DARK_TEXT, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)
    return txBox


def _add_code_block(slide, left, top, width, height, code_text, font_size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top  = Pt(8)
    tf.margin_right  = Pt(12)
    tf.margin_bottom = Pt(8)
    for i, line in enumerate(code_text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)
        p.font.name = "Consolas"
        p.space_after = Pt(2)
    return shape


def _diagram_box(slide, left, top, width, height, text,
                 fill_color=ACCENT_BLUE, text_color=WHITE, font_size=14, bold=True):
    shape = _add_shape(slide, left, top, width, height, fill_color=fill_color)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.space_before = Pt(4)
    return shape


def _arrow(slide, left, top, width, text="►", color=MID_GRAY, font_size=16):
    _add_textbox(slide, left, top, width, Inches(0.4),
                 text, font_size=font_size, color=color, alignment=PP_ALIGN.CENTER)


def _content_slide(slide, number, title):
    _set_slide_bg(slide, OFF_WHITE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BG
    bar.line.fill.background()
    _add_textbox(slide, Inches(0.8), Inches(0.15), Inches(10.5), Inches(0.8),
                 title, font_size=28, color=WHITE, bold=True)
    _add_textbox(slide, Inches(11.5), Inches(0.25), Inches(1.5), Inches(0.5),
                 f"SLIDE {number}", font_size=11, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)


def _letsdoit_slide(slide, number, exercise_id, title, tasks):
    _set_slide_bg(slide, LETSDOIT_BG)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5),
                 f"SLIDE {number}", font_size=12, color=MID_GRAY)
    _add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.0),
                 "LET'S DO IT", font_size=48, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(0.8),
                 f"{exercise_id}: {title}", font_size=24, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)
    _add_bullets(slide, Inches(3.0), Inches(4.0), Inches(7), Inches(2.5),
                 tasks, font_size=20, color=WHITE)


def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 1 — Section Header: D3-M7 Web Chat UI
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
_set_slide_bg(s, DARK_BG)

# Accent bar at top
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

_add_textbox(s, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
             "D3-M7  ·  30 minutes", font_size=14, color=MID_GRAY)

_add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.8),
             "From CLI to Web Chat UI",
             font_size=48, color=WHITE, bold=True)

_add_textbox(s, Inches(0.8), Inches(3.5), Inches(11), Inches(0.8),
             "Take the completed PV Agent and expose it through a browser-based\nchat interface — with real-time streaming via Server-Sent Events (SSE).",
             font_size=20, color=LIGHT_GRAY)

# Learning objectives strip
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.4), SLIDE_W, Inches(1.8))
strip.fill.solid()
strip.fill.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x1A)
strip.line.fill.background()

_add_textbox(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5),
             "You will learn: Switch SDK Console → Web  ·  Per-session state with ConcurrentDictionary  ·  Stream tokens via SSE",
             font_size=16, color=ACCENT_TEAL)


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 2 — Architecture: CLI vs Web
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
_content_slide(s, 2, "Architecture: One SDK Change, a New Surface")

# ── LEFT SIDE: CLI (lab 06) ──
_add_textbox(s, Inches(0.8), Inches(1.3), Inches(4.5), Inches(0.5),
             "Lab 06 — Console App", font_size=18, color=MID_GRAY, bold=True)

for i, (label, clr) in enumerate([
    ("Microsoft.NET.Sdk\n(Console)", ACCENT_TEAL),
    ("while(true) { input() }", MID_GRAY),
    ("agent.RunStreamingAsync", ACCENT_BLUE),
    ("Console.Write each token", MID_GRAY),
]):
    _diagram_box(s, Inches(0.5), Inches(1.9 + i * 1.15), Inches(4.5), Inches(1.0),
                 label, fill_color=clr, font_size=14)
    if i < 3:
        _arrow(s, Inches(0.5), Inches(2.85 + i * 1.15), Inches(4.5), "▼")

# ── ARROW between ──
_add_textbox(s, Inches(5.2), Inches(3.4), Inches(2.5), Inches(0.6),
             "⟹ upgrade", font_size=22, color=ACCENT_ORANGE, bold=True,
             alignment=PP_ALIGN.CENTER)

# ── RIGHT SIDE: Web App ──
_add_textbox(s, Inches(8.0), Inches(1.3), Inches(4.8), Inches(0.5),
             "Lab 07 — Web App", font_size=18, color=ACCENT_BLUE, bold=True)

for i, (label, clr) in enumerate([
    ("Microsoft.NET.Sdk.Web\n(only change in .csproj!)", ACCENT_BLUE),
    ("POST /chat  +  Static Files", ACCENT_TEAL),
    ("agent.RunStreamingAsync", ACCENT_BLUE),
    ("SSE data: <token>\\n\\n", ACCENT_ORANGE),
]):
    _diagram_box(s, Inches(8.0), Inches(1.9 + i * 1.15), Inches(5.0), Inches(1.0),
                 label, fill_color=clr, font_size=14)
    if i < 3:
        _arrow(s, Inches(8.0), Inches(2.85 + i * 1.15), Inches(5.0), "▼")

_add_textbox(s, Inches(0.5), Inches(6.5), Inches(12.5), Inches(0.6),
             "The agent logic (instructions + tools) is UNCHANGED. Only the transport layer switches from stdin/stdout to HTTP.",
             font_size=16, color=DARK_TEXT, bold=True)


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 3 — TODO 1 & 2: Agent Creation + Session Store
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
_content_slide(s, 3, "TODO 1 & 2 — Create the Agent & Session Store")

# ── Annotation badge ──
badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(0.5), Inches(1.2), Inches(3.2), Inches(0.55))
badge.fill.solid()
badge.fill.fore_color.rgb = ACCENT_ORANGE
badge.line.fill.background()
_add_textbox(s, Inches(0.5), Inches(1.25), Inches(3.2), Inches(0.5),
             "Why .AsIChatClient() first?", font_size=13, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

_add_textbox(s, Inches(0.5), Inches(1.85), Inches(5.5), Inches(1.0),
             "ASP.NET Core adds Microsoft.Extensions.AI, introducing\nextra AsAIAgent() overloads. Inserting .AsIChatClient()\nresolves the type ambiguity explicitly.",
             font_size=13, color=DARK_TEXT)

# ── Code: TODO 1 ──
_add_textbox(s, Inches(0.5), Inches(2.95), Inches(6), Inches(0.4),
             "TODO 1 — Create AIAgent:", font_size=14, color=DARK_TEXT, bold=True)
_add_code_block(s, Inches(0.5), Inches(3.4), Inches(6.2), Inches(2.3),
                'AIAgent agent = new OpenAIClient(\n'
                '    new ApiKeyCredential(apiKey),\n'
                '    new OpenAIClientOptions { Endpoint = new Uri(endpoint) })\n'
                '  .GetChatClient(deploymentName)\n'
                '  .AsIChatClient()          // ← resolves overload ambiguity\n'
                '  .AsAIAgent(\n'
                '      instructions: pvAgentInstructions,\n'
                '      name: "PVAgent",\n'
                '      tools: [\n'
                '          AIFunctionFactory.Create(GetProjectBudget),\n'
                '          AIFunctionFactory.Create(SubmitPv)\n'
                '      ]);',
                font_size=12)

# ── Right: TODO 2 ──
_add_textbox(s, Inches(7.0), Inches(1.2), Inches(6.0), Inches(0.4),
             "TODO 2 — Session Store:", font_size=14, color=DARK_TEXT, bold=True)
_add_code_block(s, Inches(7.0), Inches(1.65), Inches(6.0), Inches(0.85),
                'var sessions =\n'
                '    new ConcurrentDictionary<string, AgentSession>();',
                font_size=14)

# ── Why ConcurrentDictionary? ──
_add_textbox(s, Inches(7.0), Inches(2.65), Inches(6.0), Inches(0.4),
             "Why ConcurrentDictionary?", font_size=14, color=DARK_TEXT, bold=True)
_add_bullets(s, Inches(7.0), Inches(3.1), Inches(6.0), Inches(2.8), [
    "Thread-safe: multiple HTTP requests may arrive simultaneously",
    "Key = sessionId sent by the browser (stored in localStorage)",
    "Value = AgentSession holding full conversation history",
    "New tab → new sessionId → fresh AgentSession",
    "Page refresh → same sessionId → session survives (memory resumes)",
], font_size=14, color=DARK_TEXT)

# ── Bottom note ──
_add_textbox(s, Inches(0.5), Inches(6.4), Inches(12.5), Inches(0.6),
             "The CLI used a single AgentSession for one user. The web version isolates sessions so multiple tabs never share history.",
             font_size=15, color=DARK_TEXT, bold=True)


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 4 — TODO 3: POST /chat  SSE Streaming
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
_content_slide(s, 4, "TODO 3 — POST /chat: Stream Tokens as Server-Sent Events")

# ── Left: code ──
_add_textbox(s, Inches(0.5), Inches(1.2), Inches(7.5), Inches(0.4),
             "Minimal endpoint — 5 steps:", font_size=14, color=DARK_TEXT, bold=True)

_add_code_block(s, Inches(0.5), Inches(1.65), Inches(7.3), Inches(5.1),
                'app.MapPost("/chat", async (HttpContext ctx) =>\n'
                '{\n'
                '    // a) Parse JSON body: { "message", "sessionId" }\n'
                '    using var doc = JsonDocument.Parse(await bodyText);\n'
                '    var message   = doc.RootElement.GetProperty("message").GetString();\n'
                '    var sessionId = doc.RootElement.GetProperty("sessionId").GetString();\n'
                '\n'
                '    // b) Look up or create session\n'
                '    if (!sessions.TryGetValue(sessionId, out var session))\n'
                '    {\n'
                '        session = await agent.CreateSessionAsync();\n'
                '        sessions[sessionId] = session;\n'
                '    }\n'
                '\n'
                '    // c) SSE headers\n'
                '    ctx.Response.ContentType = "text/event-stream";\n'
                '    ctx.Response.Headers["Cache-Control"] = "no-cache";\n'
                '\n'
                '    // d) Stream tokens\n'
                '    await foreach (var update in\n'
                '        agent.RunStreamingAsync(message, session))\n'
                '    {\n'
                '        var data = JsonSerializer.Serialize(update.Text);\n'
                '        await ctx.Response.WriteAsync($"data: {data}\\n\\n");\n'
                '        await ctx.Response.Body.FlushAsync(); // send immediately\n'
                '    }\n'
                '\n'
                '    // e) Signal end-of-stream\n'
                '    await ctx.Response.WriteAsync("data: [DONE]\\n\\n");\n'
                '});',
                font_size=11)

# ── Right: SSE flow diagram ──
_add_textbox(s, Inches(8.2), Inches(1.2), Inches(4.8), Inches(0.4),
             "SSE token-by-token flow:", font_size=14, color=DARK_TEXT, bold=True)

flow_items = [
    ("Browser\nPOST /chat", ACCENT_BLUE),
    ("Parse body\nFind / create session", ACCENT_TEAL),
    ("RunStreamingAsync", ACCENT_BLUE),
    ('data: "Hello"\\n\\n\ndata: " there"\\n\\n', ACCENT_ORANGE),
    ("data: [DONE]\\n\\n", RGBColor(0x10, 0x7C, 0x10)),
]
for i, (label, clr) in enumerate(flow_items):
    _diagram_box(s, Inches(8.2), Inches(1.7 + i * 1.05), Inches(4.8), Inches(0.9),
                 label, fill_color=clr, font_size=12)
    if i < len(flow_items) - 1:
        _arrow(s, Inches(8.2), Inches(2.55 + i * 1.05), Inches(4.8), "▼", font_size=12)

_add_textbox(s, Inches(0.5), Inches(6.5), Inches(12.5), Inches(0.6),
             "FlushAsync() after every token ensures the browser receives and renders each word as it arrives — not all at once.",
             font_size=15, color=DARK_TEXT, bold=True)


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 5 — Let's Do It: Run, Test, Verify
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
_letsdoit_slide(s, 5, "D3-M7", "Integrate PV Agent into Web Chat UI", [
    "Copy appsettings.json.example → appsettings.json and fill credentials",
    "Implement TODO 1: create AIAgent with .AsIChatClient().AsAIAgent(...)",
    "Implement TODO 2: create ConcurrentDictionary session store",
    "Implement TODO 3: map POST /chat with SSE streaming",
    "Run dotnet run → open http://localhost:5000 → submit a full PV",
    "Verify new document in Cosmos DB Data Explorer",
    "Open a second tab — confirm independent conversation history",
])


# ═══════════════════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "webapp_slides.pptx")
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
