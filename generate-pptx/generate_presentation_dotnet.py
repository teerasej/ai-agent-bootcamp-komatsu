#!/usr/bin/env python3
"""
Generate the PV Agent & MA Agent Workshop Presentation — C# / .NET Edition (49 slides).
Mirrors generate_presentation.py but with C# code examples and dotnet lab structure.
The final section (slides 45–49) is the Web Chat UI content from generate_webapp_slides.py.

Run: python3 generate_presentation_dotnet.py
Output: presentation_dotnet.pptx in the same directory.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── colour palette ────────────────────────────────────────────────
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE     = RGBColor(0xF5, 0xF5, 0xF5)
DARK_BG       = RGBColor(0x1E, 0x1E, 0x2E)
ACCENT_BLUE   = RGBColor(0x00, 0x78, 0xD4)
ACCENT_TEAL   = RGBColor(0x00, 0xB4, 0xA0)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
LIGHT_GRAY    = RGBColor(0xE0, 0xE0, 0xE0)
MID_GRAY      = RGBColor(0x99, 0x99, 0x99)
DARK_TEXT     = RGBColor(0x22, 0x22, 0x22)
CODE_BG       = RGBColor(0x2D, 0x2D, 0x2D)
SECTION_BG    = RGBColor(0x00, 0x3E, 0x6B)
LETSDOIT_BG   = RGBColor(0x00, 0x5A, 0x9E)
TRACK1_CLR    = RGBColor(0x00, 0x78, 0xD4)
TRACK2_CLR    = RGBColor(0x10, 0x7C, 0x10)
COSMOS_CLR    = RGBColor(0x6B, 0x2F, 0xA0)
DOTNET_PURPLE = RGBColor(0x51, 0x2B, 0xD4)   # .NET brand purple accent

# ── slide dimensions (widescreen 16:9) ───────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H


# ═══════════════════════════════════════════════════════════════════
#  Helpers
# ═══════════════════════════════════════════════════════════════════

def _set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, fill_color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullets(slide, left, top, width, height, items, font_size=16,
                 color=DARK_TEXT, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)
    return txBox


def _add_code_block(slide, left, top, width, height, code_text, font_size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left   = Pt(12)
    tf.margin_top    = Pt(8)
    tf.margin_right  = Pt(12)
    tf.margin_bottom = Pt(8)
    for i, line in enumerate(code_text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)
        p.font.name = "Consolas"
        p.space_after = Pt(2)
    return shape


def _section_header(slide, number, title, subtitle="", bg_color=SECTION_BG):
    _set_slide_bg(slide, bg_color)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 f"SLIDE {number}", font_size=12, color=MID_GRAY)
    _add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(1.5),
                 title, font_size=40, color=WHITE, bold=True)
    if subtitle:
        _add_textbox(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(1.0),
                     subtitle, font_size=20, color=LIGHT_GRAY)


def _content_slide(slide, number, title, bg_color=OFF_WHITE):
    _set_slide_bg(slide, bg_color)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BG
    bar.line.fill.background()
    _add_textbox(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.8),
                 title, font_size=28, color=WHITE, bold=True)
    _add_textbox(slide, Inches(11.5), Inches(0.25), Inches(1.5), Inches(0.5),
                 f"SLIDE {number}", font_size=11, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


def _letsdoit_slide(slide, number, exercise_id, title, tasks, bg_color=LETSDOIT_BG):
    _set_slide_bg(slide, bg_color)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5),
                 f"SLIDE {number}", font_size=12, color=MID_GRAY)
    _add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.0),
                 "LET'S DO IT", font_size=48, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(0.8),
                 f"{exercise_id}: {title}", font_size=24, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)
    _add_bullets(slide, Inches(3.0), Inches(4.0), Inches(7), Inches(2.5),
                 tasks, font_size=20, color=WHITE)


def _diagram_box(slide, left, top, width, height, text,
                 fill_color=ACCENT_BLUE, text_color=WHITE, font_size=14, bold=True):
    shape = _add_shape(slide, left, top, width, height, fill_color=fill_color)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.space_before = Pt(4)
    return shape


def _arrow_text(slide, left, top, width, text, color=MID_GRAY, font_size=12):
    _add_textbox(slide, left, top, width, Inches(0.4),
                 text, font_size=font_size, color=color, alignment=PP_ALIGN.CENTER)


def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])   # blank layout


# ═══════════════════════════════════════════════════════════════════
#  SLIDES
# ═══════════════════════════════════════════════════════════════════

# ──────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ──────────────────────────────────────────────────────────────────
s = new_slide()
_set_slide_bg(s, DARK_BG)

# thin .NET purple top bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
bar.fill.solid()
bar.fill.fore_color.rgb = DOTNET_PURPLE
bar.line.fill.background()

_add_textbox(s, Inches(0.8), Inches(0.5), Inches(4), Inches(0.5),
             "C# / .NET Edition", font_size=14, color=DOTNET_PURPLE, bold=True)

_add_textbox(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5),
             "Building AI Agents for\nPV Creation & Approval",
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
_add_textbox(s, Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.7),
             "Hands-on Workshop",
             font_size=24, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)
_add_textbox(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.5),
             "AI Agent Bootcamp — Komatsu  ·  Microsoft Agent Framework for .NET",
             font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────────────────────────
# SLIDE 2 — The Business Problem
# ──────────────────────────────────────────────────────────────────
s = new_slide()
_content_slide(s, 2, "The Business Problem")
_add_textbox(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.8),
             "Payment Voucher requests are manual, error-prone, and slow to approve.",
             font_size=26, color=DARK_TEXT, bold=True)
_add_bullets(s, Inches(1.2), Inches(2.8), Inches(10), Inches(3.5), [
    "Manual data entry — requestors fill forms by hand, missing fields are common",
    "No validation — incomplete or incorrect PVs reach the manager's desk",
    "Slow approval chain — managers manually search through submissions",
    "No single source of truth — data scattered across emails and spreadsheets",
], font_size=18, color=DARK_TEXT)

# ──────────────────────────────────────────────────────────────────
# SLIDE 3 — The Solution: Two AI Agents
# ──────────────────────────────────────────────────────────────────
s = new_slide()
_content_slide(s, 3, "The Solution: Two AI Agents")
_diagram_box(s, Inches(1.5), Inches(2.0), Inches(3.0), Inches(1.5),
             "PV Agent\n(Requestor Side)", fill_color=TRACK1_CLR)
_arrow_text(s, Inches(1.5), Inches(3.6), Inches(3.0), "creates PV  ──►")
_diagram_box(s, Inches(5.2), Inches(2.0), Inches(3.0), Inches(1.5),
             "Azure\nCosmos DB", fill_color=COSMOS_CLR)
_diagram_box(s, Inches(8.8), Inches(2.0), Inches(3.0), Inches(1.5),
             "MA Agent\n(Manager Side)", fill_color=TRACK2_CLR)
_arrow_text(s, Inches(8.8), Inches(3.6), Inches(3.0), "◄──  reviews & approves")
_add_textbox(s, Inches(1.5), Inches(4.5), Inches(3.0), Inches(0.5),
             "Requestor", font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
_add_textbox(s, Inches(8.8), Inches(4.5), Inches(3.0), Inches(0.5),
             "Manager", font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
_add_textbox(s, Inches(0.8), Inches(5.5), Inches(11), Inches(0.6),
             "We'll build two agents — one for creation, one for review. They share a database.",
             font_size=18, color=DARK_TEXT, bold=True)

# ══════════════════════════════════════════════════════════════════
#  TRACK 1: PV AGENT
# ══════════════════════════════════════════════════════════════════

# ──────────────────────────────────────────────────────────────────
# SLIDE 4 — What is an AI Agent?
# ──────────────────────────────────────────────────────────────────
s = new_slide()
_content_slide(s, 4, "What is an AI Agent?")
_add_textbox(s, Inches(0.8), Inches(1.8), Inches(11), Inches(0.8),
             "AI Agent  =  LLM  +  Instructions  +  Tools",
             font_size=32, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
for i, (label, color) in enumerate([
    ("LLM\n(the brain)", ACCENT_BLUE),
    ("Instructions\n(the rules)", ACCENT_TEAL),
    ("Tools\n(the hands)", ACCENT_ORANGE),
]):
    _diagram_box(s, Inches(2.0 + i * 3.5), Inches(3.5), Inches(2.8), Inches(1.5),
                 label, fill_color=color, font_size=18)

# ──────────────────────────────────────────────────────────────────
# SLIDE 5 — PV Agent Boundaries
# ──────────────────────────────────────────────────────────────────
s = new_slide()
_content_slide(s, 5, "PV Agent Boundaries")
boxes = [
    ("User (requestor)", MID_GRAY),
    ("PV Agent  (Microsoft Agent Framework + Azure AI Foundry)", TRACK1_CLR),
    ("Validation Layer", ACCENT_TEAL),
    ("PV System / Back-office (mock)", MID_GRAY),
]
for i, (label, color) in enumerate(boxes):
    _diagram_box(s, Inches(3.0), Inches(1.5 + i * 1.2), Inches(7.0), Inches(0.9),
                 label, fill_color=color, font_size=15)
    if i < len(boxes) - 1:
        _arrow_text(s, Inches(3.0), Inches(2.35 + i * 1.2), Inches(7.0), "▼", font_size=16)
_add_textbox(s, Inches(0.8), Inches(6.2), Inches(11), Inches(0.6),
             "The agent is a data assistant. It does NOT approve, does NOT access finance systems.",
             font_size=18, color=DARK_TEXT, bold=True)

# ──────────────────────────────────────────────────────────────────
# SLIDE 6 — Let's Do It: Setup
# ──────────────────────────────────────────────────────────────────
s = new_slide()
_letsdoit_slide(s, 6, "D3-M1", "Environment Setup (C#)", [
    "Create a Microsoft Foundry project",
    "Deploy GPT-4.1 model",
    "Configure appsettings.json",
    "dotnet build  →  dotnet run  (smoke test)",
])

# ──────────────────────────────────────────────────────────────────
# Exercise 2 — Agent Foundation (Slides 7–12)
# ──────────────────────────────────────────────────────────────────

# SLIDE 7 — The System Prompt is Everything
s = new_slide()
_section_header(s, 7, "The System Prompt\nis Everything",
                "It defines who the agent is, what it knows, and what it refuses to do.")

# SLIDE 8 — Anatomy of a Good Instruction
s = new_slide()
_content_slide(s, 8, "Anatomy of a Good Instruction")
layers = [
    ("1. Role", '"You are PV Agent..."', ACCENT_BLUE),
    ("2. Required Fields", "List every field the agent must collect", ACCENT_TEAL),
    ("3. Constraints", "What the agent must NOT do", ACCENT_ORANGE),
    ("4. Output Format", "The exact JSON template", COSMOS_CLR),
]
for i, (label, desc, color) in enumerate(layers):
    _diagram_box(s, Inches(1.5), Inches(1.5 + i * 1.3), Inches(3.5), Inches(1.0),
                 label, fill_color=color, font_size=18)
    _add_textbox(s, Inches(5.5), Inches(1.6 + i * 1.3), Inches(7), Inches(0.8),
                 desc, font_size=18, color=DARK_TEXT)

# SLIDE 9 — Slot-Filling Pattern
s = new_slide()
_content_slide(s, 9, "Slot-Filling Conversation Pattern")
steps = [
    "1. Greet &\nunderstand intent",
    "2. Identify PV type\n(MonthlyFee / OneTime)",
    "3. Collect missing fields\n— ask ONE at a time",
    "4. Show\nconfirmation summary",
    "5. Output\nfinal JSON",
]
for i, step in enumerate(steps):
    _diagram_box(s, Inches(0.7 + i * 2.4), Inches(2.5), Inches(2.2), Inches(1.8),
                 step, fill_color=ACCENT_BLUE if i < 4 else ACCENT_TEAL, font_size=12)
    if i < len(steps) - 1:
        _arrow_text(s, Inches(2.7 + i * 2.4), Inches(3.2), Inches(0.5), "►", font_size=16)
_add_textbox(s, Inches(0.8), Inches(5.5), Inches(11), Inches(0.6),
             "The agent follows a slot-filling pattern — collect one field at a time, never invent data.",
             font_size=18, color=DARK_TEXT, bold=True)

# SLIDE 10 — Design Rules
s = new_slide()
_content_slide(s, 10, "Instruction Design Rules")
_add_bullets(s, Inches(1.2), Inches(1.8), Inches(10), Inches(4.5), [
    "✓  Never ask for a field the user already provided",
    "✓  Never invent data that the user did not provide",
    '✓  State assumptions explicitly  (e.g. "I\'ll use THB as the default currency")',
    "✓  Only output the final JSON after user confirms the summary",
], font_size=22, color=DARK_TEXT)

# SLIDE 11 — Agent Framework SDK Pattern (C#)
s = new_slide()
_content_slide(s, 11, "Microsoft Agent Framework SDK — Minimal Pattern (C#)")
_add_code_block(s, Inches(1.5), Inches(1.5), Inches(10), Inches(4.2),
                '// Create the AIAgent\n'
                'AIAgent agent = new OpenAIClient(\n'
                '    new ApiKeyCredential(apiKey),\n'
                '    new OpenAIClientOptions { Endpoint = new Uri(endpoint) })\n'
                '    .GetChatClient(deploymentName)\n'
                '    .AsAIAgent(\n'
                '        instructions: pvAgentInstructions,\n'
                '        name: "PVAgent",\n'
                '        tools: [...]);\n'
                '\n'
                '// AgentSession holds the full conversation history\n'
                'AgentSession session = await agent.CreateSessionAsync();\n'
                '\n'
                '// Conversation loop — stream tokens\n'
                'await foreach (AgentResponseUpdate update in\n'
                '    agent.RunStreamingAsync(userInput, session))\n'
                '{\n'
                '    Console.Write(update.Text);\n'
                '}',
                font_size=15)
_add_textbox(s, Inches(1.5), Inches(6.0), Inches(10), Inches(0.6),
             "Three ingredients: OpenAIClient connection, instruction string, streaming conversation loop.",
             font_size=18, color=DARK_TEXT, bold=True)

# SLIDE 12 — Let's Do It: Agent Foundation
s = new_slide()
_letsdoit_slide(s, 12, "D3-M2", "Agent Role & Conversation Design (C#)", [
    "Write the PV Agent instruction (system prompt) as a C# raw string literal",
    "Test 3 user stories in Foundry playground",
    "Implement the multi-turn loop with AIAgent + AgentSession",
])

# ──────────────────────────────────────────────────────────────────
# Exercise 3 — Data Grounding (Slides 13–16)
# ──────────────────────────────────────────────────────────────────

# SLIDE 13 — New Requirement: Budget Type
s = new_slide()
_content_slide(s, 13, "New Requirement: Budget Type")
_add_textbox(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.8),
             "Finance says: every PV now needs a budgetType — either Expense or Investment.",
             font_size=24, color=DARK_TEXT, bold=True)
_diagram_box(s, Inches(1.5), Inches(3.0), Inches(4.5), Inches(2.5),
             "Expense\n\nSubscriptions, fees,\ntraining registrations,\nconsumable items",
             fill_color=ACCENT_TEAL, font_size=16)
_diagram_box(s, Inches(7.0), Inches(3.0), Inches(4.5), Inches(2.5),
             "Investment\n\nEquipment purchases,\ntool licenses,\ninfrastructure upgrades",
             fill_color=ACCENT_ORANGE, font_size=16)

# SLIDE 14 — Instruction Grounding vs. Code Logic
s = new_slide()
_content_slide(s, 14, "Instruction Grounding vs. Code Logic")
_diagram_box(s, Inches(1.0), Inches(1.8), Inches(5.0), Inches(1.0),
             "Code approach (C# if/else)", fill_color=MID_GRAY, font_size=18)
_add_bullets(s, Inches(1.2), Inches(3.0), Inches(5.0), Inches(2.0), [
    'if (desc.Contains("monthly")) type = "Expense";',
    "Misses edge cases",
    '"one-time monitoring upgrade" → ???',
], font_size=16, color=DARK_TEXT)
_diagram_box(s, Inches(7.0), Inches(1.8), Inches(5.0), Inches(1.0),
             "Instruction approach", fill_color=ACCENT_TEAL, font_size=18)
_add_bullets(s, Inches(7.2), Inches(3.0), Inches(5.0), Inches(2.0), [
    "Add rule to the prompt (raw string literal)",
    "LLM understands nuance",
    "Handles ambiguity naturally",
], font_size=16, color=DARK_TEXT)
_add_textbox(s, Inches(0.8), Inches(5.8), Inches(11), Inches(0.6),
             "If the decision depends on understanding meaning, put it in the instruction — not code.",
             font_size=18, color=DARK_TEXT, bold=True)

# SLIDE 15 — What Changes: Instruction Only
s = new_slide()
_content_slide(s, 15, "What Changes: Instruction Only")
_add_textbox(s, Inches(0.8), Inches(1.8), Inches(11), Inches(0.8),
             "Zero code changes. Only the instruction string is updated.",
             font_size=28, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)
_add_textbox(s, Inches(0.8), Inches(3.0), Inches(11), Inches(0.6),
             "REQUIRED FIELDS section of the raw string literal — add:", font_size=18, color=DARK_TEXT, bold=True)
_add_code_block(s, Inches(1.2), Inches(3.6), Inches(10.5), Inches(1.5),
                '- expense.budgetType: MUST be "Expense" or "Investment"\n'
                '    - "Expense" for recurring/operational costs\n'
                '    - "Investment" for one-time capital expenditures',
                font_size=14)
_add_textbox(s, Inches(0.8), Inches(5.3), Inches(11), Inches(0.6),
             "The cheapest, fastest way to extend agent output — no recompile needed.",
             font_size=18, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

# SLIDE 16 — Let's Do It: Data Grounding
s = new_slide()
_letsdoit_slide(s, 16, "D3-M3", "Data Grounding via Instruction (C#)", [
    "Add budgetType classification rule to the C# raw string literal",
    "Update OUTPUT FORMAT JSON template in the instruction",
    "Test with 3 user stories (Expense, Investment, Ambiguous)",
])

# ──────────────────────────────────────────────────────────────────
# Exercise 4 — Function Tool: Read (Slides 17–21)
# ──────────────────────────────────────────────────────────────────

# SLIDE 17 — When Instructions Aren't Enough
s = new_slide()
_content_slide(s, 17, "When Instructions Aren't Enough")
_add_textbox(s, Inches(0.8), Inches(1.8), Inches(11), Inches(0.8),
             "Instructions handle rules.\nBut the agent can't look up facts from thin air.",
             font_size=26, color=DARK_TEXT, bold=True)
_add_textbox(s, Inches(1.2), Inches(3.5), Inches(10), Inches(1.5),
             'Example:\n'
             'User: "Project is IT Internal"\n'
             "Agent needs: budget = 100,000 THB, remaining = 100,000 THB\n"
             "→ The agent doesn't know this. It needs a data source.",
             font_size=18, color=DARK_TEXT)

# SLIDE 18 — Decision Matrix
s = new_slide()
_content_slide(s, 18, "Instruction vs. Tool — Decision Matrix")
_diagram_box(s, Inches(1.5), Inches(2.0), Inches(4.5), Inches(1.5),
             "Classification rules, enums\n\n→ Instruction grounding",
             fill_color=ACCENT_TEAL, font_size=18)
_diagram_box(s, Inches(7.0), Inches(2.0), Inches(4.5), Inches(1.5),
             "Data lookup, calculations, I/O\n\n→ Function tool (C# method)",
             fill_color=ACCENT_ORANGE, font_size=18)

# SLIDE 19 — The [Description] Attribute Pattern (C#)
s = new_slide()
_content_slide(s, 19, "[Description] Attribute Pattern (C#)")
_add_code_block(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(3.7),
                '[Description("Look up the total budget and remaining budget\n'
                '              for a given project from the CSV data file.")]\n'
                'static string GetProjectBudget(\n'
                '    [Description("The name of the project to look up")] string projectName)\n'
                '{\n'
                '    // Read CSV from AppContext.BaseDirectory/data/\n'
                '    string[] lines = File.ReadAllLines(dataPath);\n'
                '    // ... find row, return JSON string\n'
                '    return $"{{\"totalBudget\": {budget}, \"remainingBudget\": {remaining}}}";\n'
                '}\n'
                '\n'
                '// Register with the agent:\n'
                'tools: [ AIFunctionFactory.Create(GetProjectBudget) ]',
                font_size=14)
_add_bullets(s, Inches(1.2), Inches(5.3), Inches(10), Inches(1.5), [
    "① [Description] on the method — tells the agent what the tool does",
    "② [Description] on each parameter — tells the agent what to pass",
    "③ Return string — the agent reads the JSON result as conversation context",
], font_size=16, color=DARK_TEXT)

# SLIDE 20 — Data Flow: Agent Calls Tool
s = new_slide()
_content_slide(s, 20, "Data Flow: Agent Calls Tool")
flow = [
    ('User says:\n"Project is IT Internal"', MID_GRAY),
    ("Agent calls\nGetProjectBudget(\n\"IT Internal\")", TRACK1_CLR),
    ("Method reads CSV\n→ returns JSON string", ACCENT_TEAL),
    ("Agent uses result\nin PV JSON output", ACCENT_BLUE),
]
for i, (label, color) in enumerate(flow):
    _diagram_box(s, Inches(0.5 + i * 3.1), Inches(2.3), Inches(2.8), Inches(2.2),
                 label, fill_color=color, font_size=13)
    if i < len(flow) - 1:
        _arrow_text(s, Inches(3.1 + i * 3.1), Inches(3.2), Inches(0.5), "►", font_size=18)
_add_textbox(s, Inches(0.8), Inches(5.5), Inches(11), Inches(0.6),
             "The agent autonomously decides when to call the tool. You don't trigger it manually.",
             font_size=18, color=DARK_TEXT, bold=True)

# SLIDE 21 — Let's Do It: Budget Lookup Tool
s = new_slide()
_letsdoit_slide(s, 21, "D3-M4", "Function Tool for Budget Lookup (C#)", [
    "Implement GetProjectBudget with [Description] attributes",
    "Register with AIFunctionFactory.Create(GetProjectBudget)",
    "Test with known & unknown project names",
])

# ──────────────────────────────────────────────────────────────────
# Exercise 5 — Function Tool: Write (Slides 22–25)
# ──────────────────────────────────────────────────────────────────

# SLIDE 22 — Reversed Data Flow
s = new_slide()
_content_slide(s, 22, "Reversed Data Flow")
_diagram_box(s, Inches(1.0), Inches(2.0), Inches(5.0), Inches(1.3),
             "Exercise 4 (Read)\nTool → returns data → Agent",
             fill_color=ACCENT_TEAL, font_size=18)
_diagram_box(s, Inches(7.0), Inches(2.0), Inches(5.0), Inches(1.3),
             "Exercise 5 (Write)\nAgent → passes JSON → SubmitPv()",
             fill_color=ACCENT_ORANGE, font_size=18)
_add_textbox(s, Inches(0.8), Inches(4.2), Inches(11), Inches(0.8),
             "Now the agent is the sender.\nIt assembles the complete PV payload and hands it to SubmitPv().",
             font_size=20, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

# SLIDE 23 — The Completeness Gate
s = new_slide()
_content_slide(s, 23, "The Completeness Gate")
_add_textbox(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "Instruction rules that prevent premature submission:", font_size=20, color=DARK_TEXT)
_add_bullets(s, Inches(1.2), Inches(2.5), Inches(10), Inches(3.5), [
    '1.  Do NOT call SubmitPv until the user confirms AND all fields are filled',
    '2.  status must be "ReadyForSubmission" before calling the tool',
    "3.  Never submit with placeholders, empty strings, or zero amounts",
    '4.  Never submit a PV with status = "Draft"',
], font_size=20, color=DARK_TEXT)
_add_textbox(s, Inches(0.8), Inches(5.5), Inches(11), Inches(0.6),
             "The instruction acts as a gate — preventing premature submission.",
             font_size=18, color=DARK_TEXT, bold=True)

# SLIDE 24 — Stub First, Wire Later
s = new_slide()
_content_slide(s, 24, "Stub First, Wire Later")
_add_textbox(s, Inches(0.8), Inches(1.8), Inches(11), Inches(0.8),
             "SubmitPv just writes JSON to the console. No database yet.",
             font_size=24, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
_add_code_block(s, Inches(1.5), Inches(2.8), Inches(10), Inches(2.2),
                '[Description("Submit PV data to system for approval")]\n'
                'async Task<string> SubmitPv(\n'
                '    [Description("Complete PV JSON as a string")] string pvJson,\n'
                '    CancellationToken ct = default)\n'
                '{\n'
                '    Console.WriteLine("--- PV Submission Received ---");\n'
                '    Console.WriteLine(pvJson);\n'
                '    Console.WriteLine("--- End of PV Submission ---");\n'
                '    return "PV submission received and ready for database insertion.";\n'
                '}',
                font_size=14)
_add_textbox(s, Inches(0.8), Inches(5.5), Inches(11), Inches(0.6),
             "Validate the payload shape before adding infrastructure complexity.",
             font_size=18, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

# SLIDE 25 — Let's Do It: Submit PV Tool
s = new_slide()
_letsdoit_slide(s, 25, "D3-M5", "Submit PV Tool (C#)", [
    "Implement SubmitPv with [Description] attributes",
    "Register both tools: AIFunctionFactory.Create(GetAllProjectBudgets), AIFunctionFactory.Create(SubmitPv)",
    "Add completeness gate rules to the instruction",
    "Test full flow: collect → confirm → submit → console output",
])

# ──────────────────────────────────────────────────────────────────
# Exercise 6 — Cosmos DB (Slides 26–29)
# ──────────────────────────────────────────────────────────────────

# SLIDE 26 — From Stub to Real Database
s = new_slide()
_content_slide(s, 26, "From Stub to Real Database")
_diagram_box(s, Inches(1.0), Inches(2.0), Inches(4.5), Inches(1.5),
             "Before\nSubmitPv → Console.WriteLine\n(terminal only)",
             fill_color=MID_GRAY, font_size=16)
_arrow_text(s, Inches(5.5), Inches(2.4), Inches(2.0), "⟹", font_size=32)
_diagram_box(s, Inches(7.5), Inches(2.0), Inches(4.5), Inches(1.5),
             "After\nSubmitPv → CosmosClient\n→ Cosmos DB",
             fill_color=COSMOS_CLR, font_size=16)
_add_textbox(s, Inches(0.8), Inches(4.5), Inches(11), Inches(0.6),
             "Same tool interface. Only the method body changes. The agent doesn't notice the swap.",
             font_size=20, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

# SLIDE 27 — Why Cosmos DB Serverless?
s = new_slide()
_content_slide(s, 27, "Why Cosmos DB Serverless?")
_add_bullets(s, Inches(1.2), Inches(1.8), Inches(10), Inches(4.0), [
    "Pay-per-request — no minimum throughput charge",
    "NoSQL — stores JSON natively (matches PV structure perfectly)",
    "Shared container — both PV Agent and MA Agent use the same data",
    "Partition key: /id — each document is its own partition (UUID via Guid.NewGuid())",
], font_size=22, color=DARK_TEXT)

# SLIDE 28 — Separation of Concerns
s = new_slide()
_content_slide(s, 28, "Separation of Concerns")
layers_28 = [
    ("Agent Instruction", "unchanged", ACCENT_TEAL),
    ("Tool Interface", "unchanged — SubmitPv(string pvJson) -> Task<string>", ACCENT_BLUE),
    ("Tool Implementation", "CHANGED — CosmosClient insert + stream API", COSMOS_CLR),
]
for i, (label, desc, color) in enumerate(layers_28):
    _diagram_box(s, Inches(1.5), Inches(1.8 + i * 1.6), Inches(4.0), Inches(1.2),
                 label, fill_color=color, font_size=18)
    _add_textbox(s, Inches(6.0), Inches(2.0 + i * 1.6), Inches(6.5), Inches(0.8),
                 desc, font_size=18, color=DARK_TEXT)
_add_textbox(s, Inches(0.8), Inches(6.2), Inches(11), Inches(0.6),
             "Clean separation. The agent's behavior is stable even as the backend evolves.",
             font_size=18, color=DARK_TEXT, bold=True)

# SLIDE 29 — Let's Do It: Cosmos DB
s = new_slide()
_letsdoit_slide(s, 29, "D3-M6", "Store PV in Cosmos DB (C#)", [
    "Provision Cosmos DB Serverless account in the Azure Portal",
    "Replace SubmitPv body with CosmosClient + CreateItemStreamAsync",
    "Verify the document in Cosmos DB Data Explorer",
])

# ──────────────────────────────────────────────────────────────────
# SLIDE 30 — PV Agent Recap
# ──────────────────────────────────────────────────────────────────
s = new_slide()
_section_header(s, 30, "PV Agent: What We Built",
                "6 layers, incrementally added — one concept per exercise.")
concepts = [
    ("Setup", ACCENT_BLUE),
    ("Instruction", ACCENT_TEAL),
    ("Data Grounding", ACCENT_ORANGE),
    ("Read Tool", TRACK1_CLR),
    ("Write Tool", COSMOS_CLR),
    ("Cosmos DB", COSMOS_CLR),
]
for i, (label, color) in enumerate(concepts):
    _diagram_box(s, Inches(0.5 + i * 2.1), Inches(5.0), Inches(1.9), Inches(1.0),
                 f"{i+1}. {label}", fill_color=color, font_size=13, text_color=WHITE)
    if i < len(concepts) - 1:
        _arrow_text(s, Inches(2.2 + i * 2.1), Inches(5.2), Inches(0.5), "►", color=WHITE, font_size=16)

# ══════════════════════════════════════════════════════════════════
#  TRACK 2: MA AGENT
# ══════════════════════════════════════════════════════════════════

# SLIDE 31 — Two Sides of the Same Workflow
s = new_slide()
_content_slide(s, 31, "Two Sides of the Same Workflow")
_diagram_box(s, Inches(1.5), Inches(2.2), Inches(3.0), Inches(1.5),
             "PV Agent\n(completed ✓)", fill_color=LIGHT_GRAY, text_color=MID_GRAY, font_size=16)
_diagram_box(s, Inches(5.2), Inches(2.2), Inches(3.0), Inches(1.5),
             "Azure\nCosmos DB", fill_color=COSMOS_CLR)
_diagram_box(s, Inches(8.8), Inches(2.2), Inches(3.0), Inches(1.5),
             "MA Agent\n(building now)", fill_color=TRACK2_CLR, font_size=16)
_add_textbox(s, Inches(0.8), Inches(5.0), Inches(11), Inches(0.6),
             "Now we build the other side.",
             font_size=24, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

# SLIDE 32 — Same Patterns, Different Role
s = new_slide()
_content_slide(s, 32, "Same Patterns, Different Role")
headers = ["", "PV Agent", "MA Agent"]
rows = [
    ["Role", "Collects & creates PV", "Reviews & approves PV"],
    ["Data direction", "Writes to Cosmos DB", "Reads from Cosmos DB"],
    ["SDK pattern", "Same (AIAgent + AgentSession)", "Same (AIAgent + AgentSession)"],
    ["Instruction", "Different", "Different"],
]
y_start = Inches(1.8)
for j, h in enumerate(headers):
    _diagram_box(s, Inches(1.0 + j * 3.8), y_start, Inches(3.5), Inches(0.7),
                 h, fill_color=DARK_BG, font_size=16, text_color=WHITE)
for i, row in enumerate(rows):
    for j, cell in enumerate(row):
        color = LIGHT_GRAY if i % 2 == 0 else OFF_WHITE
        _diagram_box(s, Inches(1.0 + j * 3.8), y_start + Inches(0.8 + i * 0.8),
                     Inches(3.5), Inches(0.7), cell, fill_color=color,
                     text_color=DARK_TEXT, font_size=13, bold=False)
_add_textbox(s, Inches(0.8), Inches(6.2), Inches(11), Inches(0.6),
             "Agent specialization = same framework, different instructions.",
             font_size=18, color=DARK_TEXT, bold=True)

# ──────────────────────────────────────────────────────────────────
# Exercise 7 — MA Agent Foundation (Slides 33–35)
# ──────────────────────────────────────────────────────────────────

# SLIDE 33 — The Manager's Instruction
s = new_slide()
_content_slide(s, 33, "The Manager's Instruction")
_add_textbox(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "Key differences from PV Agent:", font_size=20, color=DARK_TEXT, bold=True)
_add_bullets(s, Inches(1.2), Inches(2.3), Inches(10), Inches(3.5), [
    'Role: "Help managers review and explore PV requests"  (not create)',
    "Understands the PV JSON structure  (reads the data contract)",
    'Constraint: "Do NOT modify, do NOT invent data"',
    'Constraint: "When manager asks to see PVs, ALWAYS call GetPvRequests tool"',
    "Focus: Summarize, flag issues, answer questions about PV content",
], font_size=20, color=DARK_TEXT)

# SLIDE 34 — Reusing the SDK Pattern (C#)
s = new_slide()
_content_slide(s, 34, "Reusing the SDK Pattern (C#)")
_add_code_block(s, Inches(1.5), Inches(1.5), Inches(10), Inches(3.5),
                '// Only the instruction and tools list differ from PV Agent\n'
                'AIAgent agent = new OpenAIClient(\n'
                '    new ApiKeyCredential(apiKey),\n'
                '    new OpenAIClientOptions { Endpoint = new Uri(endpoint) })\n'
                '    .GetChatClient(deploymentName)\n'
                '    .AsAIAgent(\n'
                '        instructions: maAgentInstructions,   // ← only this changes\n'
                '        name: "MAAgent",\n'
                '        tools: [...]);                        // ← added next\n'
                '\n'
                '// Same session + streaming loop as PV Agent',
                font_size=15)
_add_textbox(s, Inches(1.5), Inches(5.3), Inches(10), Inches(0.6),
             "Identical code structure. Only the instruction text changes.",
             font_size=20, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

# SLIDE 35 — Let's Do It: MA Agent Foundation
s = new_slide()
_letsdoit_slide(s, 35, "D4-M1", "MA Agent Setup & Foundation (C#)", [
    "Write the MA Agent instruction raw string literal for the reviewer role",
    "Implement conversation loop with AIAgent + AgentSession",
    "Test with manager review scenarios (paste PV JSON, ask questions)",
])

# ──────────────────────────────────────────────────────────────────
# Exercise 8 — MA Agent Function Tools (Slides 36–39)
# ──────────────────────────────────────────────────────────────────

# SLIDE 36 — A Manager Needs Two Actions
s = new_slide()
_content_slide(s, 36, "A Manager Needs Two Actions")
_diagram_box(s, Inches(1.0), Inches(2.0), Inches(5.0), Inches(2.5),
             'Query\n\n"Show me all pending PV requests"\n\n→ GetPvRequests(status)',
             fill_color=ACCENT_TEAL, font_size=16)
_diagram_box(s, Inches(7.0), Inches(2.0), Inches(5.0), Inches(2.5),
             'Act\n\n"Approve PV pv-001"\n\n→ UpdatePvApprovalStatus(id, status)',
             fill_color=ACCENT_ORANGE, font_size=16)
_add_textbox(s, Inches(0.8), Inches(5.5), Inches(11), Inches(0.6),
             "One agent, two tools — read and write.",
             font_size=20, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

# SLIDE 37 — Prototype with In-Memory Data (C#)
s = new_slide()
_content_slide(s, 37, "Prototype with In-Memory Data (C#)")
_add_textbox(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "Before connecting to Cosmos DB, use a C# List<string> as the data source.",
             font_size=20, color=DARK_TEXT, bold=True)
_add_code_block(s, Inches(0.8), Inches(2.3), Inches(11.5), Inches(3.2),
                '// Mutable list — each element is a raw JSON string (one PV document)\n'
                'List<string> samplePvData =\n'
                '[\n'
                '    """{"id":"pv-001","pvTitle":"Monthly GitHub Copilot Subscription",\n'
                '         "approval":{"status":"Pending"},...}""",\n'
                '    """{"id":"pv-002","pvTitle":"Conference Registration Fee",\n'
                '         "approval":{"status":"Approved"},...}""",\n'
                '    """{"id":"pv-003","pvTitle":"Ergonomic Office Chairs",\n'
                '         "approval":{"status":"Approved"},...}"""\n'
                '];',
                font_size=13)
_add_textbox(s, Inches(0.8), Inches(6.0), Inches(11), Inches(0.6),
             "Prototype fast. Validate tool patterns and agent behavior without database complexity.",
             font_size=18, color=DARK_TEXT, bold=True)

# SLIDE 38 — Registering Multiple Tools (C#)
s = new_slide()
_content_slide(s, 38, "Registering Multiple Tools (C#)")
_add_code_block(s, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                'tools: [\n'
                '    AIFunctionFactory.Create(GetPvRequests),\n'
                '    AIFunctionFactory.Create(UpdatePvApprovalStatus)\n'
                ']',
                font_size=18)
_add_textbox(s, Inches(0.8), Inches(4.2), Inches(11), Inches(1.2),
             'The agent decides which tool to call based on the user\'s question.\n'
             '"Show me pending PVs"  →  GetPvRequests\n'
             '"Approve pv-001"  →  UpdatePvApprovalStatus',
             font_size=18, color=DARK_TEXT)

# SLIDE 39 — Let's Do It: MA Agent Tools
s = new_slide()
_letsdoit_slide(s, 39, "D4-M2", "Function Tools for PV Review & Approval (C#)", [
    "Implement GetPvRequests and UpdatePvApprovalStatus with [Description] attributes",
    "Register both tools with AIFunctionFactory.Create",
    "Test: query pending → approve → verify approval",
])

# ──────────────────────────────────────────────────────────────────
# Exercise 9 — MA Agent + Cosmos DB (Slides 40–42)
# ──────────────────────────────────────────────────────────────────

# SLIDE 40 — The Full Pipeline
s = new_slide()
_content_slide(s, 40, "The Full Pipeline — End to End")
_diagram_box(s, Inches(0.4), Inches(1.8), Inches(2.5), Inches(1.0),
             "PV Agent", fill_color=TRACK1_CLR, font_size=16)
_diagram_box(s, Inches(0.4), Inches(3.2), Inches(2.5), Inches(0.8),
             "SubmitPv()", fill_color=ACCENT_TEAL, font_size=14)
_arrow_text(s, Inches(0.4), Inches(2.8), Inches(2.5), "▼", font_size=14)
_diagram_box(s, Inches(3.8), Inches(1.9), Inches(5.5), Inches(2.6),
             "Azure Cosmos DB\n(payment-vouchers)\n\nCreateItemStreamAsync ← PV Agent\nQueryItemsAsync → MA Agent\nReplaceItemAsync ← MA Agent",
             fill_color=COSMOS_CLR, font_size=13)
_diagram_box(s, Inches(10.3), Inches(1.8), Inches(2.8), Inches(1.0),
             "MA Agent", fill_color=TRACK2_CLR, font_size=16)
_diagram_box(s, Inches(10.3), Inches(3.2), Inches(2.8), Inches(0.8),
             "GetPvRequests\nUpdateApproval", fill_color=ACCENT_ORANGE, font_size=12)
_arrow_text(s, Inches(10.3), Inches(2.8), Inches(2.8), "▼", font_size=14)
_add_textbox(s, Inches(0.8), Inches(5.5), Inches(11), Inches(0.6),
             "Both agents now share live data. Submit a PV, then review and approve it.",
             font_size=18, color=DARK_TEXT, bold=True)

# SLIDE 41 — What Changed: Implementation Only
s = new_slide()
_content_slide(s, 41, "What Changed: Implementation Only")
layers_41 = [
    ("Agent Instruction", "unchanged", ACCENT_TEAL),
    ("Tool Interface", "unchanged — same [Description] signatures", ACCENT_BLUE),
    ("Tool Implementation", "CHANGED — List<string> → CosmosClient queries", COSMOS_CLR),
]
for i, (label, desc, color) in enumerate(layers_41):
    _diagram_box(s, Inches(1.5), Inches(1.8 + i * 1.6), Inches(4.0), Inches(1.2),
                 label, fill_color=color, font_size=18)
    _add_textbox(s, Inches(6.0), Inches(2.0 + i * 1.6), Inches(6.5), Inches(0.8),
                 desc, font_size=18, color=DARK_TEXT)
_add_textbox(s, Inches(0.8), Inches(6.2), Inches(11), Inches(0.6),
             "Same pattern as PV Agent's Cosmos upgrade. The swap is always in the implementation layer.",
             font_size=18, color=DARK_TEXT, bold=True)

# SLIDE 42 — Let's Do It: MA Agent + Cosmos DB
s = new_slide()
_letsdoit_slide(s, 42, "D4-M3", "Connect MA Agent to Cosmos DB (C#)", [
    "Add GetCosmosContainer() helper that returns the Container client",
    "Replace both tools with Cosmos DB QueryItemsAsync / PatchItemAsync",
    "End-to-end: PV Agent submit → MA Agent query → approve → verify",
])

# ──────────────────────────────────────────────────────────────────
# SLIDE 43 — What We Built Today
# ──────────────────────────────────────────────────────────────────
s = new_slide()
_content_slide(s, 43, "What We Built Today")
_diagram_box(s, Inches(1.0), Inches(1.8), Inches(3.0), Inches(1.2),
             "PV Agent\n6 exercises", fill_color=TRACK1_CLR, font_size=16)
_diagram_box(s, Inches(5.0), Inches(1.8), Inches(3.0), Inches(1.2),
             "Cosmos DB\nshared container", fill_color=COSMOS_CLR, font_size=16)
_diagram_box(s, Inches(9.0), Inches(1.8), Inches(3.0), Inches(1.2),
             "MA Agent\n3 exercises", fill_color=TRACK2_CLR, font_size=16)
concepts_all = [
    "System Prompt Design",
    "Instruction Grounding",
    "Function Tools (Read)",
    "Function Tools (Write)",
    "Cloud Persistence",
    "Agent Specialization",
]
for i, concept in enumerate(concepts_all):
    col = i % 3
    row = i // 3
    _diagram_box(s, Inches(1.5 + col * 3.6), Inches(3.8 + row * 1.2),
                 Inches(3.2), Inches(0.9),
                 f"{i+1}. {concept}",
                 fill_color=ACCENT_BLUE if i < 5 else TRACK2_CLR, font_size=14)

# SLIDE 44 — Key Takeaways
s = new_slide()
_section_header(s, 44, "Key Takeaways", bg_color=DARK_BG)
_add_bullets(s, Inches(1.2), Inches(4.0), Inches(10), Inches(3.0), [
    "1.  The system prompt is your most powerful tool",
    "2.  Classification → instruction grounding.  Facts → C# method with [Description].",
    "3.  Stub first (Console.WriteLine), wire the database later (CosmosClient)",
    "4.  Same AIAgent SDK + different instruction = different agent personality",
    "5.  Agents specialize — separate creation from review",
], font_size=22, color=WHITE, bold=True)

# ══════════════════════════════════════════════════════════════════
#  SECTION 3: WEB CHAT UI  (from generate_webapp_slides.py)
#  Slides 45–49
# ══════════════════════════════════════════════════════════════════

# ──────────────────────────────────────────────────────────────────
# SLIDE 45 — Section Header: D3-M7 Web Chat UI
# ──────────────────────────────────────────────────────────────────
s = new_slide()
_set_slide_bg(s, DARK_BG)

bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

_add_textbox(s, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
             "D3-M7  ·  30 minutes", font_size=14, color=MID_GRAY)
_add_textbox(s, Inches(11.5), Inches(0.4), Inches(1.5), Inches(0.5),
             "SLIDE 45", font_size=11, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

_add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.8),
             "From CLI to Web Chat UI",
             font_size=48, color=WHITE, bold=True)
_add_textbox(s, Inches(0.8), Inches(3.5), Inches(11), Inches(0.8),
             "Take the completed PV Agent and expose it through a browser-based\n"
             "chat interface — with real-time streaming via Server-Sent Events (SSE).",
             font_size=20, color=LIGHT_GRAY)

strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.4), SLIDE_W, Inches(1.8))
strip.fill.solid()
strip.fill.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x1A)
strip.line.fill.background()
_add_textbox(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5),
             "You will learn:  Switch SDK Console → Web  ·  Per-session state with ConcurrentDictionary  ·  Stream tokens via SSE",
             font_size=16, color=ACCENT_TEAL)

# ──────────────────────────────────────────────────────────────────
# SLIDE 46 — Architecture: CLI vs Web
# ──────────────────────────────────────────────────────────────────
s = new_slide()
_content_slide(s, 46, "Architecture: One SDK Change, a New Surface")

_add_textbox(s, Inches(0.8), Inches(1.3), Inches(4.5), Inches(0.5),
             "Lab 06 — Console App", font_size=18, color=MID_GRAY, bold=True)
for i, (label, clr) in enumerate([
    ("Microsoft.NET.Sdk\n(Console)", ACCENT_TEAL),
    ("while(true) { Console.ReadLine() }", MID_GRAY),
    ("agent.RunStreamingAsync", ACCENT_BLUE),
    ("Console.Write(update.Text)", MID_GRAY),
]):
    _diagram_box(s, Inches(0.5), Inches(1.9 + i * 1.15), Inches(4.5), Inches(1.0),
                 label, fill_color=clr, font_size=13)
    if i < 3:
        _arrow_text(s, Inches(0.5), Inches(2.85 + i * 1.15), Inches(4.5), "▼")

_add_textbox(s, Inches(5.2), Inches(3.4), Inches(2.5), Inches(0.6),
             "⟹ upgrade", font_size=22, color=ACCENT_ORANGE, bold=True,
             alignment=PP_ALIGN.CENTER)

_add_textbox(s, Inches(8.0), Inches(1.3), Inches(4.8), Inches(0.5),
             "Lab 07 — Web App", font_size=18, color=ACCENT_BLUE, bold=True)
for i, (label, clr) in enumerate([
    ("Microsoft.NET.Sdk.Web\n(only change in .csproj!)", ACCENT_BLUE),
    ("POST /chat  +  Static Files", ACCENT_TEAL),
    ("agent.RunStreamingAsync", ACCENT_BLUE),
    ('SSE  data: "token"\\n\\n', ACCENT_ORANGE),
]):
    _diagram_box(s, Inches(8.0), Inches(1.9 + i * 1.15), Inches(5.0), Inches(1.0),
                 label, fill_color=clr, font_size=13)
    if i < 3:
        _arrow_text(s, Inches(8.0), Inches(2.85 + i * 1.15), Inches(5.0), "▼")

_add_textbox(s, Inches(0.5), Inches(6.5), Inches(12.5), Inches(0.6),
             "The agent logic (instructions + tools) is UNCHANGED. Only the transport layer switches from stdin/stdout to HTTP.",
             font_size=16, color=DARK_TEXT, bold=True)

# ──────────────────────────────────────────────────────────────────
# SLIDE 47 — TODO 1 & 2: Agent Creation + Session Store
# ──────────────────────────────────────────────────────────────────
s = new_slide()
_content_slide(s, 47, "TODO 1 & 2 — Create the Agent & Session Store")

# Badge
badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(0.5), Inches(1.2), Inches(3.5), Inches(0.55))
badge.fill.solid()
badge.fill.fore_color.rgb = ACCENT_ORANGE
badge.line.fill.background()
_add_textbox(s, Inches(0.5), Inches(1.25), Inches(3.5), Inches(0.5),
             "Why .AsIChatClient() first?", font_size=13, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

_add_textbox(s, Inches(0.5), Inches(1.85), Inches(6.0), Inches(1.0),
             "ASP.NET Core adds Microsoft.Extensions.AI, introducing extra\n"
             "AsAIAgent() overloads. Inserting .AsIChatClient() resolves\n"
             "the type ambiguity explicitly.",
             font_size=13, color=DARK_TEXT)

_add_textbox(s, Inches(0.5), Inches(2.95), Inches(6), Inches(0.4),
             "TODO 1 — Create AIAgent:", font_size=14, color=DARK_TEXT, bold=True)
_add_code_block(s, Inches(0.5), Inches(3.4), Inches(6.2), Inches(2.3),
                'AIAgent agent = new OpenAIClient(\n'
                '    new ApiKeyCredential(apiKey),\n'
                '    new OpenAIClientOptions { Endpoint = new Uri(endpoint) })\n'
                '  .GetChatClient(deploymentName)\n'
                '  .AsIChatClient()          // ← resolves overload ambiguity\n'
                '  .AsAIAgent(\n'
                '      instructions: pvAgentInstructions,\n'
                '      name: "PVAgent",\n'
                '      tools: [\n'
                '          AIFunctionFactory.Create(GetProjectBudget),\n'
                '          AIFunctionFactory.Create(SubmitPv)\n'
                '      ]);',
                font_size=12)

_add_textbox(s, Inches(7.0), Inches(1.2), Inches(6.0), Inches(0.4),
             "TODO 2 — Session Store:", font_size=14, color=DARK_TEXT, bold=True)
_add_code_block(s, Inches(7.0), Inches(1.65), Inches(6.0), Inches(0.9),
                'var sessions =\n'
                '    new ConcurrentDictionary<string, AgentSession>();',
                font_size=14)

_add_textbox(s, Inches(7.0), Inches(2.65), Inches(6.0), Inches(0.4),
             "Why ConcurrentDictionary?", font_size=14, color=DARK_TEXT, bold=True)
_add_bullets(s, Inches(7.0), Inches(3.1), Inches(6.0), Inches(2.8), [
    "Thread-safe: multiple HTTP requests arrive simultaneously",
    "Key = sessionId from browser (localStorage)",
    "Value = AgentSession with full conversation history",
    "New tab → new sessionId → fresh AgentSession",
    "Page refresh → same sessionId → session survives",
], font_size=14, color=DARK_TEXT)

_add_textbox(s, Inches(0.5), Inches(6.4), Inches(12.5), Inches(0.6),
             "The CLI used a single AgentSession for one user. The web version isolates sessions so multiple tabs never share history.",
             font_size=15, color=DARK_TEXT, bold=True)

# ──────────────────────────────────────────────────────────────────
# SLIDE 48 — TODO 3: POST /chat  SSE Streaming
# ──────────────────────────────────────────────────────────────────
s = new_slide()
_content_slide(s, 48, "TODO 3 — POST /chat: Stream Tokens as Server-Sent Events")

_add_textbox(s, Inches(0.5), Inches(1.2), Inches(7.5), Inches(0.4),
             "Minimal endpoint — 5 steps:", font_size=14, color=DARK_TEXT, bold=True)

_add_code_block(s, Inches(0.5), Inches(1.65), Inches(7.3), Inches(5.1),
                'app.MapPost("/chat", async (HttpContext ctx) =>\n'
                '{\n'
                '    // a) Parse JSON body: { "message", "sessionId" }\n'
                '    using var doc = JsonDocument.Parse(bodyText);\n'
                '    var message   = doc.RootElement.GetProperty("message").GetString();\n'
                '    var sessionId = doc.RootElement.GetProperty("sessionId").GetString();\n'
                '\n'
                '    // b) Look up or create session\n'
                '    if (!sessions.TryGetValue(sessionId, out var session))\n'
                '    {\n'
                '        session = await agent.CreateSessionAsync();\n'
                '        sessions[sessionId] = session;\n'
                '    }\n'
                '\n'
                '    // c) SSE headers\n'
                '    ctx.Response.ContentType = "text/event-stream";\n'
                '    ctx.Response.Headers["Cache-Control"] = "no-cache";\n'
                '\n'
                '    // d) Stream tokens\n'
                '    await foreach (var update in\n'
                '        agent.RunStreamingAsync(message, session))\n'
                '    {\n'
                '        var data = JsonSerializer.Serialize(update.Text);\n'
                '        await ctx.Response.WriteAsync($"data: {data}\\n\\n");\n'
                '        await ctx.Response.Body.FlushAsync();\n'
                '    }\n'
                '\n'
                '    // e) Signal end-of-stream\n'
                '    await ctx.Response.WriteAsync("data: [DONE]\\n\\n");\n'
                '});',
                font_size=11)

_add_textbox(s, Inches(8.2), Inches(1.2), Inches(4.8), Inches(0.4),
             "SSE token-by-token flow:", font_size=14, color=DARK_TEXT, bold=True)

flow_items = [
    ("Browser\nPOST /chat", ACCENT_BLUE),
    ("Parse body\nFind / create session", ACCENT_TEAL),
    ("RunStreamingAsync", ACCENT_BLUE),
    ('data: "Hello"\\n\\n\ndata: " there"\\n\\n', ACCENT_ORANGE),
    ("data: [DONE]\\n\\n", RGBColor(0x10, 0x7C, 0x10)),
]
for i, (label, clr) in enumerate(flow_items):
    _diagram_box(s, Inches(8.2), Inches(1.7 + i * 1.05), Inches(4.8), Inches(0.9),
                 label, fill_color=clr, font_size=12)
    if i < len(flow_items) - 1:
        _arrow_text(s, Inches(8.2), Inches(2.55 + i * 1.05), Inches(4.8), "▼", font_size=12)

_add_textbox(s, Inches(0.5), Inches(6.5), Inches(12.5), Inches(0.6),
             "FlushAsync() after every token ensures the browser receives and renders each word as it arrives — not all at once.",
             font_size=15, color=DARK_TEXT, bold=True)

# ──────────────────────────────────────────────────────────────────
# SLIDE 49 — Let's Do It: Web App
# ──────────────────────────────────────────────────────────────────
s = new_slide()
_letsdoit_slide(s, 49, "D3-M7", "Integrate PV Agent into Web Chat UI", [
    "Copy appsettings.json.example → appsettings.json and fill credentials",
    "Implement TODO 1: create AIAgent with .AsIChatClient().AsAIAgent(...)",
    "Implement TODO 2: create ConcurrentDictionary session store",
    "Implement TODO 3: map POST /chat with SSE streaming",
    "dotnet run → open http://localhost:5000 → submit a full PV",
    "Verify new document in Cosmos DB Data Explorer",
    "Open a second tab — confirm independent conversation history",
])


# ═══════════════════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "presentation_dotnet.pptx")
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
